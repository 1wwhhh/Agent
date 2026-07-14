from __future__ import annotations

import pytest

pytest.skip("PPT template indexer package is currently not part of active runtime baseline", allow_module_level=True)

import json
from pathlib import Path
import shutil
from uuid import uuid4

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.ppt.template_indexer import TemplateIndexer


def _build_template(path: Path) -> None:
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]

    slide1 = presentation.slides.add_slide(blank_layout)
    box1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box1.text_frame.text = "标题：{{title}}"
    box2 = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    box2.text_frame.text = "结论：{{summary}} / {{title}}"

    slide2 = presentation.slides.add_slide(blank_layout)
    table_shape = slide2.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2))
    table_shape.table.cell(0, 0).text = "{{kpi_1}}"
    table_shape.table.cell(0, 1).text = "{{kpi_2}}"
    table_shape.table.cell(1, 0).text = "无占位符"

    slide3 = presentation.slides.add_slide(blank_layout)
    box3 = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box3.text_frame.text = "plain text"

    presentation.save(str(path))


def test_template_indexer_extracts_placeholders() -> None:
    workdir = Path("outputs") / "test_template_indexer" / uuid4().hex
    workdir.mkdir(parents=True, exist_ok=True)
    try:
        template_path = workdir / "template.pptx"
        _build_template(template_path)

        inventory = TemplateIndexer().scan(template_path)

        assert inventory.slide_count == 3
        assert inventory.placeholders["slide_1"] == ["summary", "title"]
        assert inventory.placeholders["slide_2"] == ["kpi_1", "kpi_2"]
        assert inventory.placeholders["slide_3"] == []
        assert inventory.placeholder_specs["title"]["role"] == "title"
        assert inventory.placeholder_specs["kpi_1"]["role"] == "metric"
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def test_template_indexer_loads_sidecar_placeholder_specs() -> None:
    workdir = Path("outputs") / "test_template_indexer" / uuid4().hex
    workdir.mkdir(parents=True, exist_ok=True)
    try:
        template_path = workdir / "template.pptx"
        _build_template(template_path)
        sidecar_path = template_path.with_suffix(".template.json")
        sidecar_path.write_text(
            json.dumps(
                {
                    "placeholder_specs": {
                        "title": {
                            "slide_ref": "slide_1",
                            "role": "title",
                            "max_chars": 12,
                            "instruction": "Main headline",
                        }
                    }
                },
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )

        inventory = TemplateIndexer().scan(template_path)

        assert inventory.placeholder_specs["title"]["max_chars"] == 12
        assert inventory.placeholder_specs["title"]["instruction"] == "Main headline"
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def test_template_indexer_missing_file_raises() -> None:
    with pytest.raises(FileNotFoundError):
        TemplateIndexer().scan("not_exists_template.pptx")


def test_template_indexer_rejects_non_pptx() -> None:
    workdir = Path("outputs") / "test_template_indexer" / uuid4().hex
    workdir.mkdir(parents=True, exist_ok=True)
    try:
        txt_path = workdir / "template.txt"
        txt_path.write_text("not a ppt", encoding="utf-8")

        with pytest.raises(ValueError):
            TemplateIndexer().scan(txt_path)
    finally:
        shutil.rmtree(workdir, ignore_errors=True)

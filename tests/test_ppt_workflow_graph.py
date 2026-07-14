from __future__ import annotations

import pytest

pytest.skip("PPT workflow internals are currently not part of active runtime baseline", allow_module_level=True)

import shutil
from pathlib import Path
from uuid import uuid4

from pptx import Presentation
from pptx.util import Inches

from app.graph.runtime_graph import GraphRuntimeDependencies, _RuntimeGraphBuilder
from app.router import TaskRouter
from app.router.capability import capability_from_tool
from app.schemas.task import TaskModel, TaskStatus, utc_now
from app.state import LangGraphState
from tests.support.tools import RuntimeTestTool


def _create_placeholder_template(path: Path) -> None:
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_layout)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    box.text_frame.text = "Title: {{slide_01_title}}"
    box2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    box2.text_frame.text = "Summary: {{slide_01_summary}}"
    presentation.save(str(path))


def _create_plain_layout_template(path: Path) -> None:
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_layout)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_box.text_frame.text = "Plain Title"
    body_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    body_box.text_frame.text = "Plain Body"
    presentation.save(str(path))


async def test_ppt_workflow_mapper_and_decomposer_append_render_task() -> None:
    workdir = Path("outputs") / "test_ppt_workflow_graph" / uuid4().hex
    workdir.mkdir(parents=True, exist_ok=True)
    try:
        template_path = workdir / "template.pptx"
        _create_placeholder_template(template_path)

        router = TaskRouter()
        text_tool = RuntimeTestTool(name="text_generate_tool", description="text tool", tags=["text"])
        reason_tool = RuntimeTestTool(name="llm_reason_tool", description="reason tool", tags=["reason"])
        await router.register_tools(
            [
                (text_tool, capability_from_tool(text_tool)),
                (reason_tool, capability_from_tool(reason_tool)),
            ]
        )
        builder = _RuntimeGraphBuilder(dependencies=GraphRuntimeDependencies(router=router, repair_llm_client=None))

        lang_state = LangGraphState.create(
            request_id=f"req_{uuid4().hex}",
            session_id=f"sess_{uuid4().hex}",
            user_input="请基于模板生成年终总结PPT",
            runtime_metadata={
                "delivery_format": "ppt",
                "template_path": str(template_path),
                "output_filename": "demo_output.pptx",
            },
        )
        lang_state.set_planned_tasks(
            [
                TaskModel(
                    task_id="task_1",
                    task_name="draft_content",
                    description="generate content",
                    tool="text_generate_tool",
                    input={"prompt": "生成总结"},
                    output_key="final_result",
                    depends_on=[],
                    priority=1,
                    status=TaskStatus.PENDING,
                    retry_count=0,
                    max_retry=1,
                    timeout=60,
                    created_at=utc_now(),
                ),
                TaskModel(
                    task_id="task_legacy_render",
                    task_name="legacy_render",
                    description="planner-generated legacy render task",
                    tool="ppt_render_tool",
                    input={
                        "template_path": "templates/year_end_summary_template.pptx",
                        "variables": {"final_result_text": "{{final_result.text}}"},
                    },
                    output_key="legacy_ppt_artifact",
                    depends_on=["task_1"],
                    priority=2,
                    status=TaskStatus.PENDING,
                    retry_count=0,
                    max_retry=2,
                    timeout=120,
                    created_at=utc_now(),
                ),
            ]
        )

        await builder._supervisor_impl(lang_state)
        assert lang_state.metadata["workflow_type"] == "PPT_WORKFLOW"

        await builder._template_mapper_impl(lang_state)
        assert "ppt_template_inventory" in lang_state.context.shared_data

        await builder._task_decomposer_impl(lang_state)
        task_ids = [task.task_id for task in lang_state.planned_tasks]
        assert "render_ppt_task" in task_ids
        assert "task_legacy_render" not in task_ids

        render_task = next(task for task in lang_state.planned_tasks if task.task_id == "render_ppt_task")
        assert render_task.tool == "ppt_render_tool"
        assert render_task.depends_on == ["task_1"]
        assert render_task.output_key == "ppt_artifact"
        assert render_task.input["render_mode"] == "placeholder_json"
        assert render_task.input["content_text"] == "{{final_result.text}}"

        content_task = next(task for task in lang_state.planned_tasks if task.task_id == "task_1")
        prompt_text = str(content_task.input.get("prompt") or "")
        context_text = str(content_task.input.get("context") or "")
        assert "Final Output Contract:" in prompt_text
        assert '"slide_01_title"' in prompt_text
        assert "Placeholder Generation Rules:" in context_text
        assert content_task.input["ppt_placeholder_generation"] is True
        assert int(content_task.input["max_tokens"]) >= 2200
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


async def test_ppt_workflow_decomposer_uses_auto_layout_when_template_has_no_placeholders() -> None:
    workdir = Path("outputs") / "test_ppt_workflow_graph" / uuid4().hex
    workdir.mkdir(parents=True, exist_ok=True)
    try:
        template_path = workdir / "plain_template.pptx"
        _create_plain_layout_template(template_path)

        router = TaskRouter()
        text_tool = RuntimeTestTool(name="text_generate_tool", description="text tool", tags=["text"])
        reason_tool = RuntimeTestTool(name="llm_reason_tool", description="reason tool", tags=["reason"])
        await router.register_tools(
            [
                (text_tool, capability_from_tool(text_tool)),
                (reason_tool, capability_from_tool(reason_tool)),
            ]
        )
        builder = _RuntimeGraphBuilder(dependencies=GraphRuntimeDependencies(router=router, repair_llm_client=None))

        lang_state = LangGraphState.create(
            request_id=f"req_{uuid4().hex}",
            session_id=f"sess_{uuid4().hex}",
            user_input="请基于模板生成学习总结PPT",
            runtime_metadata={
                "delivery_format": "ppt",
                "template_path": str(template_path),
                "output_filename": "auto_layout_demo.pptx",
            },
        )
        lang_state.set_planned_tasks(
            [
                TaskModel(
                    task_id="task_1",
                    task_name="draft_content",
                    description="generate content",
                    tool="text_generate_tool",
                    input={"prompt": "生成总结"},
                    output_key="final_result",
                    depends_on=[],
                    priority=1,
                    status=TaskStatus.PENDING,
                    retry_count=0,
                    max_retry=1,
                    timeout=60,
                    created_at=utc_now(),
                ),
            ]
        )

        await builder._supervisor_impl(lang_state)
        await builder._template_mapper_impl(lang_state)
        await builder._task_decomposer_impl(lang_state)

        render_task = next(task for task in lang_state.planned_tasks if task.task_id == "render_ppt_task")
        assert render_task.input["render_mode"] == "auto_layout"
        assert render_task.input["content_text"] == "{{final_result.text}}"
        assert render_task.input["variables"] == {}

        content_task = next(task for task in lang_state.planned_tasks if task.task_id == "task_1")
        context_text = str(content_task.input.get("context") or "")
        assert "PPT Template Guidance:" in context_text
        assert "fill_mode: auto_layout_fill" in context_text
    finally:
        shutil.rmtree(workdir, ignore_errors=True)

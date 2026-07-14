from __future__ import annotations

import pytest

pytest.skip("PPT render module is currently not part of active runtime baseline", allow_module_level=True)

import asyncio
import shutil
from pathlib import Path
from uuid import uuid4

from pptx import Presentation
from pptx.util import Inches

from app.tools.ppt_render import PPTRenderTool


def _create_placeholder_template(path: Path) -> None:
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]

    slide = presentation.slides.add_slide(blank_layout)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_box.text_frame.text = "Title: {{title}}"
    summary_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    summary_box.text_frame.text = "Summary: {{summary}}"

    presentation.save(str(path))


def _create_layout_template(path: Path) -> None:
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]

    slide_one = presentation.slides.add_slide(blank_layout)
    slide_one.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(1)).text_frame.text = "Cover Title"
    slide_one.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2)).text_frame.text = "Cover Summary"

    slide_two = presentation.slides.add_slide(blank_layout)
    slide_two.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(1)).text_frame.text = "Section Title"
    slide_two.shapes.add_textbox(Inches(1), Inches(2), Inches(3.5), Inches(2)).text_frame.text = "Body Left"
    slide_two.shapes.add_textbox(Inches(4.8), Inches(2), Inches(3.5), Inches(2)).text_frame.text = "Body Right"

    presentation.save(str(path))


def test_ppt_render_tool_success_replaces_placeholders() -> None:
    workdir = Path("outputs") / "test_ppt_render_tool" / uuid4().hex
    workdir.mkdir(parents=True, exist_ok=True)
    try:
        template_path = workdir / "template.pptx"
        output_dir = workdir / "out"
        _create_placeholder_template(template_path)

        tool = PPTRenderTool()
        result = asyncio.run(
            tool.arun(
                payload={
                    "template_path": str(template_path),
                    "output_dir": str(output_dir),
                    "output_filename": "final_report.pptx",
                    "variables": {"title": "Year End Summary", "summary": "Achievements and review"},
                },
                context=None,
            )
        )

        assert result.success is True
        output = result.output
        assert isinstance(output, dict)
        output_path = Path(str(output["file_path"]))
        assert output_path.exists()

        rendered = Presentation(str(output_path))
        slide_text = "\n".join(shape.text for shape in rendered.slides[0].shapes if getattr(shape, "text", ""))
        assert "Year End Summary" in slide_text
        assert "Achievements and review" in slide_text
        assert output["missing_placeholders"] == []
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def test_ppt_render_tool_can_fill_placeholders_from_json_content_text() -> None:
    workdir = Path("outputs") / "test_ppt_render_tool" / uuid4().hex
    workdir.mkdir(parents=True, exist_ok=True)
    try:
        template_path = workdir / "template_json.pptx"
        output_dir = workdir / "out"
        _create_placeholder_template(template_path)

        tool = PPTRenderTool()
        result = asyncio.run(
            tool.arun(
                payload={
                    "template_path": str(template_path),
                    "output_dir": str(output_dir),
                    "output_filename": "json_mapping_report.pptx",
                    "render_mode": "placeholder_json",
                    "content_text": '{"title": "学习总结", "summary": "阶段成果与后续计划"}',
                    "variables": {},
                },
                context=None,
            )
        )

        assert result.success is True
        output = result.output
        assert isinstance(output, dict)
        output_path = Path(str(output["file_path"]))
        rendered = Presentation(str(output_path))
        slide_text = "\n".join(shape.text for shape in rendered.slides[0].shapes if getattr(shape, "text", ""))
        assert "学习总结" in slide_text
        assert "阶段成果与后续计划" in slide_text
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def test_ppt_render_tool_reports_missing_placeholders_in_non_strict_mode() -> None:
    workdir = Path("outputs") / "test_ppt_render_tool" / uuid4().hex
    workdir.mkdir(parents=True, exist_ok=True)
    try:
        template_path = workdir / "template.pptx"
        output_dir = workdir / "out"
        _create_placeholder_template(template_path)

        tool = PPTRenderTool()
        result = asyncio.run(
            tool.arun(
                payload={
                    "template_path": str(template_path),
                    "output_dir": str(output_dir),
                    "output_filename": "final_report_missing.pptx",
                    "strict_placeholders": False,
                    "variables": {"title": "Title Only"},
                },
                context=None,
            )
        )

        assert result.success is True
        output = result.output
        assert isinstance(output, dict)
        assert output["missing_placeholders"] == ["summary"]
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def test_ppt_render_tool_fails_in_strict_mode_when_placeholder_missing() -> None:
    workdir = Path("outputs") / "test_ppt_render_tool" / uuid4().hex
    workdir.mkdir(parents=True, exist_ok=True)
    try:
        template_path = workdir / "template.pptx"
        output_dir = workdir / "out"
        _create_placeholder_template(template_path)

        tool = PPTRenderTool()
        result = asyncio.run(
            tool.arun(
                payload={
                    "template_path": str(template_path),
                    "output_dir": str(output_dir),
                    "output_filename": "strict_fail.pptx",
                    "strict_placeholders": True,
                    "variables": {"title": "Title Only"},
                },
                context=None,
            )
        )

        assert result.success is False
        assert result.error is not None
        assert "missing placeholders" in result.error
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def test_ppt_render_tool_auto_layout_fills_plain_template() -> None:
    workdir = Path("outputs") / "test_ppt_render_tool" / uuid4().hex
    workdir.mkdir(parents=True, exist_ok=True)
    try:
        template_path = workdir / "layout_template.pptx"
        output_dir = workdir / "out"
        _create_layout_template(template_path)

        tool = PPTRenderTool()
        result = asyncio.run(
            tool.arun(
                payload={
                    "template_path": str(template_path),
                    "output_dir": str(output_dir),
                    "output_filename": "auto_layout.pptx",
                    "render_mode": "auto_layout",
                    "content_text": "\n".join(
                        [
                            "---幻灯片---",
                            "第1页：学习总结",
                            "- Python 基础与项目实践",
                            "- NLP 入门与文本处理",
                            "---幻灯片---",
                            "第2页：下一步计划",
                            "- 深入学习大模型应用",
                            "- 持续完善作品集",
                        ]
                    ),
                    "variables": {},
                },
                context=None,
            )
        )

        assert result.success is True
        output = result.output
        assert isinstance(output, dict)
        output_path = Path(str(output["file_path"]))
        assert output_path.exists()

        rendered = Presentation(str(output_path))
        slide_one_text = "\n".join(shape.text for shape in rendered.slides[0].shapes if getattr(shape, "text", ""))
        slide_two_text = "\n".join(shape.text for shape in rendered.slides[1].shapes if getattr(shape, "text", ""))
        assert "学习总结" in slide_one_text
        assert "Python 基础与项目实践" in slide_one_text
        assert "下一步计划" in slide_two_text
        assert "持续完善作品集" in slide_two_text
    finally:
        shutil.rmtree(workdir, ignore_errors=True)
